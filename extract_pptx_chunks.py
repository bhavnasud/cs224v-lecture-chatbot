import json
from pptx import Presentation
import glob
import os

INPUT_DIR = "slides"
OUTPUT_DIR = "chunked_slide_transcriptions"
# Set default metadata
TITLE = "Slides"
LANGUAGE = "en"


slides = {
    "Caregiving in Cancer Med 275 FINAL.pptx": {
        "name": "Developing Culturally Attuned Interventions to Support Caregiving in Cancer",
        "last_edit_date": "2024-10-16",
        "author": "Ranak Trivedi (PhD, FSBM, FGSA)",
        "url": "https://canvas.stanford.edu/courses/198463/files?preview=14095971"
    },
    "MED 275 - Dr. Lin slides Sep 25.pptx": {
        "name": "MED 275 – Diagnosis and Screening",
        "last_edit_date": "2024-09-25",
        "author": "Bryant Lin (MD)",
        "url": "https://canvas.stanford.edu/courses/198463/files?preview=14018782"
    },
    "MED 275 - Dr. Lui slides Sep 25.pptx": {
        "name": "Lung cancer staging and screening",
        "last_edit_date": "2024-09-25",
        "author": "Natalie Liu (MD)",
        "url": "https://canvas.stanford.edu/courses/198463/files?preview=14018783"
    },
    "Stanford MED 275 Epidemiology 12.56.51 PM.pptx": {
        "name": "MED 275 - Epidemiology and Cultural Considerations",
        "last_edit_date": "2024-10-30",
        "author": "Jeffrey Velotta (MD, FACS)",
        "url": "https://canvas.stanford.edu/courses/198463/files?preview=14191559"
    },
}

    
slides_data = []
# Process each PowerPoint file in the input directory
for pptx_file in slides:
    # Load the presentation
    presentation = Presentation(os.path.join(INPUT_DIR, pptx_file))

    words_for_chunk = 0
    text_for_chunk = ""
    chunk_start_slide = 1

    # Iterate through each slide in the presentation with their index
    for slide_index, slide in enumerate(presentation.slides, start=1):  # Start counting from 1
        slide_text = []
        
        # Collect all text in the slide, sorting shapes by their vertical position
        shapes_with_text = [
            shape for shape in slide.shapes if hasattr(shape, "text_frame")
        ]
        
        # Sort shapes by their top position (y-coordinate)
        shapes_with_text.sort(key=lambda s: s.top)

        # Collect text from sorted shapes
        for shape in shapes_with_text:
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = ''.join(run.text for run in paragraph.runs)
                if paragraph_text:  # Only add non-empty paragraphs
                    slide_text.append(paragraph_text)

        # Join all text into a single string
        text = "\n".join(slide_text)
        text_for_chunk += ("\n" + text if text_for_chunk else text)
        num_words = len(text.split())
        
        words_for_chunk += num_words
        if (words_for_chunk > 300) or slide_index == len(presentation.slides):
            entry = {
                "document_title": TITLE,
                "full_section_title": f"{TITLE} > {slides[pptx_file]['name']} > Slide {chunk_start_slide}-{slide_index}",
                "content": text_for_chunk,
                "block_type": "text",
                "language": LANGUAGE,
                "last_edit_date": slides[pptx_file]["last_edit_date"],
                "url": slides[pptx_file]["url"],
                "author": slides[pptx_file]["author"]
            }
            slides_data.append(entry)
            words_for_chunk = 0
            text_for_chunk = ""
            chunk_start_slide = slide_index + 1

# Define output file path and write JSON lines format data
output_file = os.path.join(OUTPUT_DIR, "slides.jsonl")
with open(output_file, "w") as json_file:
    for entry in slides_data:
        json.dump(entry, json_file, ensure_ascii=False)
        json_file.write("\n")
